#!/usr/bin/env python3
"""
Comprehensive Document Upload Testing for ChatPDF Application
This script tests all the requirements from the review request.
"""

import requests
import json
import os
import uuid
from dotenv import load_dotenv
from reportlab.pdfgen import canvas
from docx import Document as DocxDocument
import openpyxl
import csv
from pptx import Presentation

# Load environment variables
load_dotenv('/app/frontend/.env')
load_dotenv('/app/backend/.env')

# Backend URL configuration
BACKEND_URL = os.environ.get('REACT_APP_BACKEND_URL')
INTERNAL_API_URL = "http://localhost:8001/api"
API_URL = INTERNAL_API_URL

print(f"Testing document upload functionality at: {API_URL}")

def create_test_session():
    """Create a test session"""
    url = f"{API_URL}/sessions"
    payload = {"title": "Comprehensive Document Upload Test"}
    
    response = requests.post(url, json=payload)
    if response.status_code != 200:
        print(f"Failed to create session: {response.status_code} - {response.text}")
        return None
    
    data = response.json()
    session_id = data["id"]
    print(f"✅ Created test session: {session_id}")
    return session_id

def create_comprehensive_test_files():
    """Create comprehensive test files for all supported formats"""
    files = {}
    
    # 1. Create PDF
    pdf_path = f"/app/comprehensive_test_{uuid.uuid4().hex[:8]}.pdf"
    c = canvas.Canvas(pdf_path)
    c.drawString(100, 750, "Comprehensive Test PDF Document")
    c.drawString(100, 700, "This PDF contains sample text for extraction testing.")
    c.drawString(100, 650, "Multiple lines to verify text extraction works correctly.")
    c.drawString(100, 600, "PDF format is essential for ChatPDF application.")
    c.drawString(100, 550, "Testing comprehensive document upload functionality.")
    c.save()
    files['pdf'] = pdf_path
    
    # 2. Create DOCX
    docx_path = f"/app/comprehensive_test_{uuid.uuid4().hex[:8]}.docx"
    doc = DocxDocument()
    doc.add_heading('Comprehensive Test DOCX Document', 0)
    doc.add_paragraph('This is a comprehensive test DOCX document for upload testing.')
    doc.add_paragraph('It contains multiple paragraphs to verify text extraction.')
    doc.add_paragraph('DOCX format support is important for document processing.')
    doc.add_paragraph('Testing all document upload functionality comprehensively.')
    doc.save(docx_path)
    files['docx'] = docx_path
    
    # 3. Create XLSX
    xlsx_path = f"/app/comprehensive_test_{uuid.uuid4().hex[:8]}.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Test Sheet"
    ws['A1'] = "Comprehensive Test XLSX Document"
    ws['A2'] = "Column A"
    ws['B2'] = "Column B"
    ws['C2'] = "Column C"
    ws['A3'] = "Data Row 1"
    ws['B3'] = "Value 1"
    ws['C3'] = "Value 2"
    ws['A4'] = "Data Row 2"
    ws['B4'] = "Value 3"
    ws['C4'] = "Value 4"
    ws['A5'] = "XLSX format testing"
    ws['B5'] = "Spreadsheet data"
    ws['C5'] = "Comprehensive test"
    wb.save(xlsx_path)
    files['xlsx'] = xlsx_path
    
    # 4. Create CSV
    csv_path = f"/app/comprehensive_test_{uuid.uuid4().hex[:8]}.csv"
    with open(csv_path, 'w', newline='') as csvfile:
        writer = csv.writer(csvfile)
        writer.writerow(['Column 1', 'Column 2', 'Column 3', 'Column 4'])
        writer.writerow(['Comprehensive Test CSV', 'Document', 'Upload', 'Testing'])
        writer.writerow(['Data Row 1', 'Value 1', 'Value 2', 'Value 3'])
        writer.writerow(['Data Row 2', 'Value 4', 'Value 5', 'Value 6'])
        writer.writerow(['CSV format', 'testing', 'comprehensive', 'complete'])
    files['csv'] = csv_path
    
    # 5. Create TXT
    txt_path = f"/app/comprehensive_test_{uuid.uuid4().hex[:8]}.txt"
    with open(txt_path, 'w') as txtfile:
        txtfile.write("Comprehensive Test TXT Document for Upload Testing\n")
        txtfile.write("This is a plain text file with multiple lines.\n")
        txtfile.write("Text format is the simplest document type.\n")
        txtfile.write("TXT files should be processed correctly by the system.\n")
        txtfile.write("Testing comprehensive document upload functionality.\n")
    files['txt'] = txt_path
    
    # 6. Create PPTX
    pptx_path = f"/app/comprehensive_test_{uuid.uuid4().hex[:8]}.pptx"
    prs = Presentation()
    
    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "Comprehensive Test PPTX Document"
    subtitle1.text = "Upload Testing Presentation"
    
    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    title2.text = "Content Slide"
    content2.text = "This presentation tests PPTX format support.\nMultiple slides with text content.\nPowerPoint format comprehensive testing."
    
    # Slide 3
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    title3.text = "Testing Slide"
    content3.text = "Comprehensive document upload functionality.\nAll formats should work correctly.\nText extraction verification."
    
    prs.save(pptx_path)
    files['pptx'] = pptx_path
    
    print(f"✅ Created comprehensive test files: {list(files.keys())}")
    return files

def test_1_supported_formats_endpoint():
    """Test 1: Test the supported formats endpoint returns all expected document types"""
    print("\n" + "="*80)
    print("TEST 1: SUPPORTED FORMATS ENDPOINT")
    print("="*80)
    
    url = f"{API_URL}/supported-formats"
    response = requests.get(url)
    
    print(f"Status: {response.status_code}")
    
    if response.status_code != 200:
        print(f"❌ FAIL: Expected status 200, got {response.status_code}")
        print(f"Error: {response.text}")
        return False
    
    data = response.json()
    print(f"Response: {json.dumps(data, indent=2)}")
    
    # Verify response structure
    if "supported_formats" not in data:
        print("❌ FAIL: Response missing 'supported_formats' field")
        return False
    
    if "message" not in data:
        print("❌ FAIL: Response missing 'message' field")
        return False
    
    supported_formats = data["supported_formats"]
    expected_formats = ['pdf', 'docx', 'xlsx', 'xls', 'csv', 'txt', 'pptx']
    
    for format_type in expected_formats:
        if format_type not in supported_formats:
            print(f"❌ FAIL: Format '{format_type}' not in supported formats")
            return False
    
    print(f"✅ PASS: All expected formats supported: {supported_formats}")
    return True

def test_2_universal_upload_endpoint(session_id, test_files):
    """Test 2: Test document upload for each format using universal endpoint"""
    print("\n" + "="*80)
    print("TEST 2: UNIVERSAL UPLOAD ENDPOINT FOR ALL FORMATS")
    print("="*80)
    
    results = {}
    
    for file_type, file_path in test_files.items():
        print(f"\n--- Testing {file_type.upper()} Upload ---")
        
        url = f"{API_URL}/sessions/{session_id}/upload-document"
        
        # Determine MIME type
        mime_types = {
            'pdf': 'application/pdf',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'csv': 'text/csv',
            'txt': 'text/plain',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        }
        
        with open(file_path, "rb") as file:
            files = {"file": (f"test.{file_type}", file, mime_types.get(file_type, 'application/octet-stream'))}
            response = requests.post(url, files=files)
        
        print(f"Status: {response.status_code}")
        
        if response.status_code == 200:
            data = response.json()
            print(f"✅ SUCCESS: {json.dumps(data, indent=2)}")
            
            # Verify response structure
            required_fields = ['message', 'filename', 'file_type', 'content_length']
            for field in required_fields:
                if field not in data:
                    print(f"❌ FAIL: Response missing '{field}' field")
                    results[file_type] = False
                    continue
            
            if data['file_type'] != file_type:
                print(f"❌ FAIL: Expected file_type '{file_type}', got '{data['file_type']}'")
                results[file_type] = False
                continue
            
            if data['content_length'] <= 0:
                print(f"❌ FAIL: Content length should be > 0, got {data['content_length']}")
                results[file_type] = False
                continue
            
            results[file_type] = True
        else:
            print(f"❌ FAIL: {response.text}")
            results[file_type] = False
    
    # Summary
    passed = sum(1 for success in results.values() if success)
    total = len(results)
    print(f"\n--- Universal Upload Summary ---")
    print(f"Passed: {passed}/{total}")
    
    for file_type, success in results.items():
        status = "✅ PASS" if success else "❌ FAIL"
        print(f"{file_type.upper()}: {status}")
    
    return passed == total

def test_3_backward_compatibility(session_id, pdf_path):
    """Test 3: Test backward compatibility - old PDF upload endpoint should still work"""
    print("\n" + "="*80)
    print("TEST 3: BACKWARD COMPATIBILITY - LEGACY PDF ENDPOINT")
    print("="*80)
    
    url = f"{API_URL}/sessions/{session_id}/upload-pdf"
    
    with open(pdf_path, "rb") as file:
        files = {"file": ("test.pdf", file, "application/pdf")}
        response = requests.post(url, files=files)
    
    print(f"Status: {response.status_code}")
    
    if response.status_code != 200:
        print(f"❌ FAIL: Expected status 200, got {response.status_code}")
        print(f"Error: {response.text}")
        return False
    
    data = response.json()
    print(f"✅ SUCCESS: {json.dumps(data, indent=2)}")
    
    # Verify response structure
    required_fields = ['message', 'filename', 'file_type', 'content_length']
    for field in required_fields:
        if field not in data:
            print(f"❌ FAIL: Response missing '{field}' field")
            return False
    
    if data['file_type'] != 'pdf':
        print(f"❌ FAIL: Expected file_type 'pdf', got '{data['file_type']}'")
        return False
    
    if data['message'] != 'PDF uploaded successfully':
        print(f"❌ FAIL: Expected message 'PDF uploaded successfully', got '{data['message']}'")
        return False
    
    print("✅ PASS: Legacy PDF endpoint working correctly")
    return True

def test_4_error_handling(session_id):
    """Test 4: Test error handling for unsupported file types"""
    print("\n" + "="*80)
    print("TEST 4: ERROR HANDLING - UNSUPPORTED FILE TYPES")
    print("="*80)
    
    # Test unsupported format
    unsupported_path = f"/app/test_unsupported_{uuid.uuid4().hex[:8]}.xyz"
    with open(unsupported_path, 'w') as f:
        f.write("This is an unsupported file format for testing error handling.")
    
    try:
        url = f"{API_URL}/sessions/{session_id}/upload-document"
        
        with open(unsupported_path, "rb") as file:
            files = {"file": ("test.xyz", file, "application/octet-stream")}
            response = requests.post(url, files=files)
        
        print(f"Status: {response.status_code}")
        
        if response.status_code != 400:
            print(f"❌ FAIL: Expected status 400 for unsupported format, got {response.status_code}")
            return False
        
        data = response.json()
        print(f"✅ SUCCESS: {json.dumps(data, indent=2)}")
        
        # Verify error message
        if "detail" not in data:
            print("❌ FAIL: Error response missing 'detail' field")
            return False
        
        error_message = data["detail"].lower()
        if "unsupported" not in error_message:
            print("❌ FAIL: Error message should mention 'unsupported'")
            return False
        
        print("✅ PASS: Error handling for unsupported format working correctly")
        return True
        
    finally:
        if os.path.exists(unsupported_path):
            os.remove(unsupported_path)

def test_5_session_integration(session_id):
    """Test 5: Test session integration - verify documents are stored in sessions"""
    print("\n" + "="*80)
    print("TEST 5: SESSION INTEGRATION")
    print("="*80)
    
    url = f"{API_URL}/sessions"
    response = requests.get(url)
    
    if response.status_code != 200:
        print(f"❌ FAIL: Failed to get sessions: {response.text}")
        return False
    
    sessions = response.json()
    test_session = next((s for s in sessions if s["id"] == session_id), None)
    
    if not test_session:
        print(f"❌ FAIL: Test session not found")
        return False
    
    print(f"Session data: {json.dumps(test_session, indent=2)}")
    
    # Check for new document fields
    new_fields = ['document_filename', 'document_content', 'document_type']
    for field in new_fields:
        if field not in test_session:
            print(f"❌ FAIL: Session missing new field '{field}'")
            return False
    
    # Check for backward compatibility fields
    old_fields = ['pdf_filename', 'pdf_content']
    for field in old_fields:
        if field not in test_session:
            print(f"❌ FAIL: Session missing backward compatibility field '{field}'")
            return False
    
    # Verify content is not empty
    if not test_session['document_content'] or len(test_session['document_content']) == 0:
        print("❌ FAIL: Document content is empty")
        return False
    
    print("✅ PASS: Session integration working correctly")
    print(f"   - Document filename: {test_session['document_filename']}")
    print(f"   - Document type: {test_session['document_type']}")
    print(f"   - Content length: {len(test_session['document_content'])} characters")
    return True

def test_6_text_extraction_accuracy(session_id, test_files):
    """Test 6: Test text extraction accuracy for different formats"""
    print("\n" + "="*80)
    print("TEST 6: TEXT EXTRACTION ACCURACY")
    print("="*80)
    
    results = {}
    
    for file_type, file_path in test_files.items():
        print(f"\n--- Testing text extraction for {file_type.upper()} ---")
        
        # Upload the document
        url = f"{API_URL}/sessions/{session_id}/upload-document"
        
        with open(file_path, "rb") as file:
            files = {"file": (f"test.{file_type}", file, f"application/{file_type}")}
            response = requests.post(url, files=files)
        
        if response.status_code != 200:
            print(f"❌ FAIL: Upload failed for {file_type}")
            results[file_type] = False
            continue
        
        # Get session to check extracted content
        session_url = f"{API_URL}/sessions"
        session_response = requests.get(session_url)
        sessions = session_response.json()
        
        test_session = next((s for s in sessions if s["id"] == session_id), None)
        if not test_session:
            print(f"❌ FAIL: Test session not found for {file_type}")
            results[file_type] = False
            continue
        
        extracted_content = test_session.get("document_content", "")
        if len(extracted_content) == 0:
            print(f"❌ FAIL: No content extracted from {file_type}")
            results[file_type] = False
            continue
        
        # Verify some expected text is present (case-insensitive)
        extracted_lower = extracted_content.lower()
        
        expected_texts = {
            'pdf': ["comprehensive test pdf", "sample text"],
            'docx': ["comprehensive test docx", "multiple paragraphs"],
            'txt': ["comprehensive test txt", "plain text file"],
            'csv': ["column 1", "comprehensive test csv"],
            'xlsx': ["comprehensive test xlsx", "column a"],
            'pptx': ["comprehensive test pptx", "content slide"]
        }
        
        expected = expected_texts.get(file_type, [])
        found_all = True
        
        for expected_text in expected:
            if expected_text not in extracted_lower:
                print(f"❌ FAIL: Expected text '{expected_text}' not found in extracted {file_type.upper()} content")
                found_all = False
        
        if found_all:
            print(f"✅ PASS: Text extraction working correctly for {file_type.upper()}")
            print(f"   Extracted {len(extracted_content)} characters")
            results[file_type] = True
        else:
            results[file_type] = False
    
    # Summary
    passed = sum(1 for success in results.values() if success)
    total = len(results)
    print(f"\n--- Text Extraction Summary ---")
    print(f"Passed: {passed}/{total}")
    
    for file_type, success in results.items():
        status = "✅ PASS" if success else "❌ FAIL"
        print(f"{file_type.upper()}: {status}")
    
    return passed == total

def main():
    """Main comprehensive test function"""
    print("="*80)
    print("COMPREHENSIVE DOCUMENT UPLOAD TESTING FOR CHATPDF")
    print("Testing all requirements from the review request")
    print("="*80)
    
    # Create test session
    session_id = create_test_session()
    if not session_id:
        print("❌ Cannot proceed without session")
        return False
    
    # Create comprehensive test files
    test_files = create_comprehensive_test_files()
    
    results = {}
    
    try:
        # Run all tests
        results["1_supported_formats"] = test_1_supported_formats_endpoint()
        results["2_universal_upload"] = test_2_universal_upload_endpoint(session_id, test_files)
        results["3_backward_compatibility"] = test_3_backward_compatibility(session_id, test_files['pdf'])
        results["4_error_handling"] = test_4_error_handling(session_id)
        results["5_session_integration"] = test_5_session_integration(session_id)
        results["6_text_extraction"] = test_6_text_extraction_accuracy(session_id, test_files)
        
    finally:
        # Clean up files
        for file_path in test_files.values():
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except:
                pass
        
        # Clean up session
        try:
            requests.delete(f"{API_URL}/sessions/{session_id}")
        except:
            pass
    
    # Print final summary
    print("\n" + "="*80)
    print("COMPREHENSIVE TEST SUMMARY")
    print("="*80)
    
    for test_name, success in results.items():
        status = "✅ PASS" if success else "❌ FAIL"
        print(f"{test_name}: {status}")
    
    total_tests = len(results)
    passed_tests = sum(1 for success in results.values() if success)
    success_rate = (passed_tests / total_tests) * 100 if total_tests > 0 else 0
    
    print(f"\nOverall Success Rate: {success_rate:.1f}% ({passed_tests}/{total_tests})")
    
    if success_rate == 100:
        print("🎉 EXCELLENT: ALL DOCUMENT UPLOAD FUNCTIONALITY TESTS PASSED!")
    elif success_rate >= 80:
        print("✅ GOOD: DOCUMENT UPLOAD FUNCTIONALITY IS WORKING WELL")
    elif success_rate >= 60:
        print("⚠️  WARNING: DOCUMENT UPLOAD FUNCTIONALITY HAS SOME ISSUES")
    else:
        print("❌ CRITICAL: DOCUMENT UPLOAD FUNCTIONALITY NEEDS SIGNIFICANT FIXES")
    
    return success_rate >= 80

if __name__ == "__main__":
    success = main()
    exit(0 if success else 1)